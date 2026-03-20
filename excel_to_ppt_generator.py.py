import pandas as pd
from pptx import Presentation
from pptx.util import Inches

# Load XLSB file
excel_file = "Airbus_SE.xlsb"  # Change to your file path
with pd.ExcelFile(excel_file, engine="pyxlsb") as xlsb:
    df = pd.read_excel(xlsb, sheet_name="Annual Balance Sheet")  # Change "Sheet1" if needed

# Create a PowerPoint presentation
ppt = Presentation()

# Add a title slide
slide_layout = ppt.slide_layouts[0]  # Title slide layout
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Excel Data Presentation"
subtitle.text = "Generated from Python script"

# Add a table slide
slide_layout = ppt.slide_layouts[5]  # Title + Content layout
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Data Summary"

# Define table position
rows, cols = df.shape
table = slide.shapes.add_table(rows+1, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)).table

# Add header row
for col_idx, column in enumerate(df.columns):
    table.cell(0, col_idx).text = column

# Add data to the table
for row_idx, row in df.iterrows():
    for col_idx, value in enumerate(row):
        table.cell(row_idx + 1, col_idx).text = str(value)

# Save the presentation
ppt.save("output.pptx")

print("PowerPoint file created successfully!")
