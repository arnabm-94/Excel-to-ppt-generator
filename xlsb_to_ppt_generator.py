import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import os

# Load XLSB file
excel_file = "Airbus_SE.xlsb"  # Change to your file path
if not os.path.exists(excel_file):
    raise FileNotFoundError("Excel file not found")


try:
    with pd.ExcelFile(excel_file, engine="pyxlsb") as xlsb:
        if not xlsb.sheet_names:
            raise ValueError("No sheets found in Excel file")
        df = pd.read_excel(xlsb, sheet_name=xlsb.sheet_names[0])  # sheet name is dynamic
except Exception as e:
    raise RuntimeError(f"Error reading Excel file: {e}")

if df.empty:
    raise ValueError("Excel sheet is empty")

max_rows = 20
df = df.head(max_rows)

# Create a PowerPoint presentation
ppt = Presentation()

# Add a title slide
slide_layout = ppt.slide_layouts[0]  # Title slide layout
slide = ppt.slides.add_slide(slide_layout)

if slide.shapes.title:
    slide.shapes.title.text = "Excel Data Presentation"

if len(slide.placeholders) > 1:
    slide.placeholders[1].text = "Generated from Python script"

# Add a table slide
slide_layout = ppt.slide_layouts[5]  # Title + Content layout
slide = ppt.slides.add_slide(slide_layout)
if slide.shapes.title:
    slide.shapes.title.text = "Data Summary"

# Define table position
rows, cols = df.shape

table = slide.shapes.add_table(rows+1, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)).table

# Add header row
for col_idx, column in enumerate(df.columns):
    table.cell(0, col_idx).text = column

# Add data to the table
for row_idx, row in enumerate(df.values):
    for col_idx, value in enumerate(row):
        if pd.isna(value):
            table.cell(row_idx + 1, col_idx).text = ""
        else:
            table.cell(row_idx + 1, col_idx).text = str(value)

# Save the presentation
try:
    ppt.save("output.pptx")
except PermissionError:
    raise PermissionError("Close 'output.pptx' if it is already open.")

print("PowerPoint file created successfully!")
